# -*- coding: utf-8 -*-
"""別紙（参考資料）：令和7年度決算との突合　― 第2回審議会資料に添付する4ページ ―

本編（六戸町_第2回審議会資料.pptx／37ページ）は経営戦略の推計値ベースのまま維持し、
令和7年度決算との突合結果はこの別紙で示す。原稿のレイアウトを流用して体裁を揃える。
"""
import copy
import json
import pathlib

from pptx import Presentation
from pptx.util import Inches, Pt

HERE = pathlib.Path(__file__).resolve().parent.parent
BASE_PPTX = HERE / "source" / "六戸町_第2回審議会資料_案.pptx"
OUT_PPTX = HERE / "output" / "六戸町_第2回審議会資料_別紙_R7決算突合.pptx"
M = json.loads((HERE / "data" / "metrics.json").read_text(encoding="utf-8"))
ACT = json.loads((HERE / "data" / "r7_actual.json").read_text(encoding="utf-8"))

# 原稿のうち流用するページ（この順に並べる）
#   4=章扉 / 12=9行8列の表 / 17=3行7列の表×2 / 10=左右2カラム
SOURCE_PAGES = [4, 12, 17, 10]

YKEY = ["R7", "R8", "R9", "R10", "R11", "R12", "R13", "R14", "R15", "R16"]
KUCHI = {"公共": [2563, 2545, 2527, 2508, 2490, 2472, 2452, 2432, 2412, 2391],
         "農集": [524, 520, 517, 513, 509, 506, 501, 497, 493, 489]}
SHARE = {"公共": [0.03184, 0.254673, 0.18661, 0.115417, 0.063915, 0.043817, 0.026321, 0.00904, 0.021102],
         "農集": [0.295786, 0.249802, 0.180081, 0.113113, 0.068836, 0.054322, 0.026809, 0.005093, 0.003074]}
RATES = {"現行": [1000, 120, 120, 130, 130, 140, 140, 140, 160],
         "①中間": [1250, 135, 138, 145, 148, 158, 163, 173, 193],
         "②中間": [1100, 150, 153, 160, 163, 173, 178, 188, 208],
         "④中間": [1200, 140, 143, 150, 153, 163, 168, 178, 198],
         "①最終": [1500, 150, 155, 160, 165, 175, 185, 205, 225],
         "②最終": [1200, 180, 185, 190, 195, 205, 215, 235, 255],
         "④最終": [1400, 160, 165, 170, 175, 185, 195, 215, 235]}


# ---------------------------------------------------------------- テキスト設定
def set_text(shape, text):
    tf = shape.text_frame
    lines = text.split("\n")
    paras = tf.paragraphs
    while len(paras) < len(lines):
        new = copy.deepcopy(paras[-1]._p)
        paras[-1]._p.addnext(new)
        paras = tf.paragraphs
    for i, line in enumerate(lines):
        p = paras[i]
        if not p.runs:
            src = next((q for q in paras if q.runs), None)
            if src is None:
                p.text = line
                continue
            p._p.append(copy.deepcopy(src.runs[0]._r))
        p.runs[0].text = line
        for r in p.runs[1:]:
            r._r.getparent().remove(r._r)
    for p in list(paras[len(lines):]):
        p._p.getparent().remove(p._p)


def sh(slide):
    one, many = {}, {}
    for s in slide.shapes:
        one.setdefault(s.name, s)
        many.setdefault(s.name, []).append(s)
    return one, many


def put(slide, mapping):
    one, _ = sh(slide)
    for name, val in mapping.items():
        if name not in one:
            raise KeyError(f"shape {name!r} not found")
        set_text(one[name], val)


def drop(slide, *names):
    one, _ = sh(slide)
    for n in names:
        el = one[n]._element
        el.getparent().remove(el)


# ---------------------------------------------------------------- 決算値の整理
def num(v):
    """会計表記（マイナスは△）."""
    return f"△{-v:,.0f}" if v < 0 else f"{v:,.0f}"


def act(biz, key):
    """決算実績（千円）."""
    return ACT[biz]["損益"][key] / 1000


def est(biz, key, i=0):
    return M[biz]["財政計画"][key][i]


FIG = {}
for _b in ("公共", "農集"):
    _in = act(_b, "下水道事業収益") - act(_b, "特別利益")
    _out = act(_b, "下水道事業費用") - act(_b, "特別損失")
    _mnt = _out - act(_b, "減価償却費") - act(_b, "支払利息") - act(_b, "その他営業外費用")
    FIG[_b] = {"経常収入": _in, "経常支出": _out, "汚水処理費": _mnt,
               "経常損益": _in - _out,
               "経常収支比率": _in / _out,
               "経費回収率": act(_b, "下水道使用料") / _mnt}


def gross(biz, rate, i):
    """基本使用料（年6回）＋超過使用料の算定額（千円・補正前）."""
    vol = est(biz, "有収水量", i)
    q = [round(vol * s) for s in SHARE[biz]]
    base = round(rate[0] * KUCHI[biz][i] * 6 / 1000, 1)
    over = sum(round(rate[j] * q[j] / 1000, 1) for j in range(1, 9))
    return round(base + over)


SCALE = {}
for _b in ("公共", "農集"):
    SCALE[_b] = {"補正率": act(_b, "下水道使用料") / gross(_b, RATES["現行"], 0),
                 "汚水処理費": FIG[_b]["汚水処理費"] / est(_b, "汚水処理費"),
                 "減価償却費": act(_b, "減価償却費") / est(_b, "減価償却費"),
                 "支払利息": act(_b, "支払利息") / est(_b, "支払利息"),
                 "長期前受金戻入": act(_b, "長期前受金戻入") / est(_b, "長期前受金戻入"),
                 "その他営業収益": act(_b, "その他営業収益"),
                 "その他営業外収益": (act(_b, "営業外収益") - act(_b, "他会計補助金")
                                      - act(_b, "長期前受金戻入"))}


def rate_of(p, i):
    """2か年改定の適用単価（R7＝現行、R8＝中間、R9以降＝最終）."""
    if i == 0:
        return RATES["現行"]
    return RATES[p + "中間"] if i == 1 else RATES[p + "最終"]


def sim(biz, p, year):
    """R7実績を起点に組み直した経費回収率・経常収支比率（2か年改定）.

    一般会計繰入金は経営戦略の前提（資本費相当）どおりとする。
    """
    i = YKEY.index(year)
    s = SCALE[biz]
    inc = gross(biz, rate_of(p, i), i) * s["補正率"]
    oz = est(biz, "汚水処理費", i) * s["汚水処理費"]
    cin = (inc + s["その他営業収益"] + est(biz, "補助金", i)
           + est(biz, "長期前受金戻入", i) * s["長期前受金戻入"] + s["その他営業外収益"])
    cout = (oz + est(biz, "減価償却費", i) * s["減価償却費"]
            + est(biz, "支払利息", i) * s["支払利息"])
    return inc / oz, cin / cout


# ---------------------------------------------------------------- ページの抽出
prs = Presentation(BASE_PPTX)
lst = prs.slides._sldIdLst
ids = list(lst)
RID = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
keep = {n - 1: ids[n - 1] for n in SOURCE_PAGES}
for i, el in enumerate(ids):
    if i not in keep:
        prs.part.drop_rel(el.get(RID))
        lst.remove(el)
for el in [keep[n - 1] for n in SOURCE_PAGES]:      # 原稿の並びを別紙の並びに直す
    lst.remove(el)
    lst.append(el)
sl = list(prs.slides)

# ================================================================ 表紙
put(sl[0], {
    "sec_num": "別紙", "sec_lbl": "APPENDIX", "sub_lbl": "六戸町 下水道事業",
    "main_ttl": "令和7年度決算との突合",
    "desc": "本別紙は、令和7年度決算（合計残高試算表・貸借対照表／令和8年3月31日）と、"
            "本編の収支シミュレーションが前提としている経営戦略の推計値を突合したものです。\n"
            "本編（37ページ）の数値は経営戦略の推計値のままとしています。"
            "決算実績を起点に組み直した場合の指標は、別紙2に試算として示しています。\n"
            "第2回 六戸町下水道使用料審議委員会　令和8年8月19日",
})
# 章番号の枠は「01」を前提とした大きさのため、全角2文字が収まるよう文字サイズを下げる
_one, _ = sh(sl[0])
for _p in _one["sec_num"].text_frame.paragraphs:
    for _r in _p.runs:
        _r.font.size = Pt(54)

# ================================================================ 別紙1　主要科目の突合
s = sl[1]
ROWS = [
    ("使用料収入", lambda b: (est(b, "使用料収入"), act(b, "下水道使用料"))),
    ("一般会計繰入金（補助金）", lambda b: (est(b, "補助金"), act(b, "他会計補助金"))),
    ("長期前受金戻入", lambda b: (est(b, "長期前受金戻入"), act(b, "長期前受金戻入"))),
    ("経常収入 (C)", lambda b: (est(b, "経常収入"), FIG[b]["経常収入"])),
    ("汚水処理費（維持管理費分）", lambda b: (est(b, "汚水処理費"), FIG[b]["汚水処理費"])),
    ("減価償却費", lambda b: (est(b, "減価償却費"), act(b, "減価償却費"))),
    ("支払利息", lambda b: (est(b, "支払利息"), act(b, "支払利息"))),
    ("経常支出 (D)", lambda b: (est(b, "経常支出"), FIG[b]["経常支出"])),
    ("経常損益 (C)-(D)", lambda b: (est(b, "経常損益"), FIG[b]["経常損益"])),
]
drop(s, "gl6", "th7", *[f"d{r}7" for r in range(9)])
d = {"hb": "  参考資料　別紙1　令和7年度決算と推計値の突合（主要科目）",
     "ttl": "令和7年度決算と経営戦略推計値の突合　主要科目（単位：千円）",
     "gl2": "公共下水道事業", "gl4": "農業集落排水事業",
     "th0": "科目", "th1": "R7推計", "th2": "R7決算", "th3": "乖離",
     "th4": "R7推計", "th5": "R7決算", "th6": "乖離"}
for r, (name, fn) in enumerate(ROWS):
    d[f"d{r}0"] = name
    for c, biz in enumerate(("公共", "農集")):
        e, a = fn(biz)
        d[f"d{r}{1+c*3}"] = num(e)
        d[f"d{r}{2+c*3}"] = num(a)
        d[f"d{r}{3+c*3}"] = f"{'＋' if a - e >= 0 else '△'}{abs(a-e):,.0f}"
d["note"] = (
    f"▶ 経費回収率　公共：推計34.4% → 決算 {FIG['公共']['経費回収率']*100:.1f}%　／　"
    f"農集：推計39.9% → 決算 {FIG['農集']['経費回収率']*100:.1f}%　"
    f"　経常収支比率　公共：推計98.5% → 決算 {FIG['公共']['経常収支比率']*100:.1f}%　／　"
    f"農集：推計82.0% → 決算 {FIG['農集']['経常収支比率']*100:.1f}%\n"
    "▶ 最大の乖離は公共の汚水処理費（維持管理費分）で、推計168,331千円に対し決算125,145千円（△25.7%）。"
    "令和6年度決算119,503千円からの増加は＋4.7%にとどまり、推計していた＋40.9%増は生じていません"
    "（差の大半は「その他（主に委託料）」／推計155,280千円→決算114,032千円）。\n"
    "▶ 農集は使用料収入が決算11,604千円と推計を22.5%下回る一方、一般会計繰入金が推計の約2.7倍でした。"
    "経常収支比率が112.6%となっているのはこのためです。※ 出所：R7公共／農集 合計残高試算表（令和8年3月31日）。")
put(s, d)
_one, _ = sh(s)
_one["ttl"].width = Inches(12.36)
for _c, (_l, _w) in enumerate([(0.42, 3.05), (3.55, 1.45), (5.05, 1.45), (6.55, 1.45),
                               (8.20, 1.45), (9.70, 1.45), (11.20, 1.45)]):
    _one[f"th{_c}"].left, _one[f"th{_c}"].width = Inches(_l), Inches(_w)
    for _r in range(9):
        _sp = _one[f"d{_r}{_c}"]
        _sp.left, _sp.width = Inches(_l), Inches(_w)
        _sp.top, _sp.height = Inches(2.10 + _r * 0.41), Inches(0.39)
_one["gl2"].left, _one["gl2"].width = Inches(3.55), Inches(4.45)
_one["gl4"].left, _one["gl4"].width = Inches(8.20), Inches(4.45)

# ================================================================ 別紙2　指標の試算
s = sl[2]
d = {"hb": "  参考資料　別紙2　決算を起点に組み直した場合の指標（試算）",
     "ttl": "令和7年度決算を起点に組み直した場合の指標　2か年改定・試算",
     "hcr": "経費回収率（目標：47%以上）　※（ ）内は本編の値",
     "her": "経常収支比率（目標：単年度100%以上）　※（ ）内は本編の値／一般会計繰入金は経営戦略の前提（資本費相当）",
     "gpub": "公共下水道", "gagr": "農業集落排水",
     "gpub2": "公共下水道", "gagr2": "農業集落排水"}
for pre in ("th", "th2"):
    for c, lab in enumerate(["年度", "①標準型", "②家庭軽減型", "④段階累進型",
                             "①標準型", "②家庭軽減型", "④段階累進型"]):
        d[f"{pre}{c}"] = lab
for r, y in enumerate(["R7", "R9", "R12"]):
    d[f"cr{r}0"] = d[f"er{r}0"] = y
    for c, (biz, p) in enumerate([(b, p) for b in ("公共", "農集") for p in ("①", "②", "④")], start=1):
        er, cr = sim(biz, p, y)
        key = "現行" if y == "R7" else p + "2"
        d[f"cr{r}{c}"] = f"{er*100:.1f}%（{M[biz]['経費回収率'][key][YKEY.index(y)]*100:.1f}%）"
        d[f"er{r}{c}"] = f"{cr*100:.1f}%（{M[biz]['経常収支比率'][key][YKEY.index(y)]*100:.1f}%）"
d["note"] = (
    "▶ 試算の方法：使用料収入は基本使用料の年間乗数を公共・農集とも6回（隔月請求）として再計算し、"
    f"R7実績÷R7算定値の補正率（公共{SCALE['公共']['補正率']:.3f}・農集{SCALE['農集']['補正率']:.3f}）を全年度に乗じています。"
    f"費用・長期前受金戻入はR7実績÷R7推計の比率（汚水処理費は公共{SCALE['公共']['汚水処理費']:.3f}・"
    f"農集{SCALE['農集']['汚水処理費']:.3f}）を全年度に乗じ、伸び率は経営戦略の想定を踏襲。R7列が決算実績を再現します。\n"
    "▶ 公共下水道は経費回収率がR9で56.7〜60.6%となり、目標47%はR8時点で達成する見込みとなります。"
    "農業集落排水は使用料収入の下方修正が費用の下方修正を上回るため、本編より低い水準（R9で46.6〜49.8%）となります。\n"
    "▶ 経常収支比率は一般会計繰入金の前提に左右されます。上表は経営戦略の前提（資本費相当）によるもので、"
    "農集の繰入金をR7実績水準（51,777千円）で置いた場合はR9で110〜111%となります。"
    "※ 有収水量・調定口数の実績が未受領のため暫定値です。")
put(s, d)

# ================================================================ 別紙3　確認事項の整理
s = sl[3]
bs = {b: ACT[b]["貸借"] for b in ("公共", "農集")}
put(s, {
    "hdr_txt": "参考資料　別紙3　決算突合を踏まえた確認事項の整理",
    "ttl": "令和7年度決算により決着した事項と、引き続き確認が必要な事項",
    "ch0": "決算により決着した事項",
    "cs0": "本編の確認事項　A（維持管理費）・F（年間乗数）・M・N（基準値）",
    "ph0": "✔ 決着した内容",
    "p00": "  • A 維持管理費の増加はR6→R7で＋4.7%（推計は＋40.9%）",
    "p01": "  • F 農集の基本使用料の年間乗数は×6（隔月請求）が正しい",
    "p02": "  • M・N R7実績を起点にすればR6実績補正が不要となり解消",
    "mh0": "✖ 決算からは判別できない点",
    "m00": "  • 汚水処理費が推計を43,186千円下回った理由（繰延べか推計過大か）",
    "m01": "  • 農集の一般会計繰入金が推計の約2.7倍である理由と今後の方針",
    "ch1": "引き続き確認が必要な事項",
    "cs1": "本編の確認事項　O（交付金）・P（三沢市）・Q（財政指標）・R（経営改善）",
    "ph1": "✔ 決算資料から判明した事項（確認事項Q・一部解決）",
    "p10": f"  • 企業債残高 公共{bs['公共']['企業債残高']/1000:,.0f}千円・農集{bs['農集']['企業債残高']/1000:,.0f}千円（R7年度末）",
    "p11": f"  • 現金預金 公共{bs['公共']['現金預金']/1000:,.0f}千円・農集{bs['農集']['現金預金']/1000:,.0f}千円",
    "p12": f"  • 未処分利益剰余金 公共{bs['公共']['当年度未処分利益剰余金']/1000:,.0f}千円・農集{bs['農集']['当年度未処分利益剰余金']/1000:,.0f}千円",
    "mh1": "✖ 未受領の資料",
    "m10": "  • 有収水量・調定口数の実績／一般会計繰入金の繰入基準別内訳",
    "m11": "  • 交付金の規模・対象事業費／更新投資計画／三沢市の使用料",
    "rec9": "▶ 本別紙の数値は令和7年度決算（合計残高試算表・貸借対照表／令和8年3月31日、出力 令和8年6月16日）に基づきます。"
            "本編（37ページ）の収支シミュレーションは経営戦略の推計値ベースのままとしており、"
            "決算実績を反映して組み直すかどうかは審議会でのご判断を踏まえて対応します。\n"
            "※ 突合の詳細は「R7決算突合.md」および成果物エビデンス集の「R7決算突合」「R7実績ベース試算」シートに収録しています。",
})

# ================================================================ レイアウト調整
SLIDE_W = 13.333
NOTE_NAMES = ("note", "rec9", "desc")
BG = ("bg1", "bg2", "hdr_bg", "hdr_line", "hdr_accent", "hb", "hl", "ha", "right_bg", "sec_panel")


def _vlen(t):
    return sum(0.5 if ord(c) < 0x2000 else 1.0 for c in t)


def _fontpt(shape, default=11.0):
    for para in shape.text_frame.paragraphs:
        for r in para.runs:
            if r.font.size:
                return r.font.size.pt
    return default


def fit_note(slide, shape):
    pt = _fontpt(shape)
    cap = max(1.0, (shape.width.inches - 0.35) * 72 / pt)
    lines = sum(max(1, -(-int(_vlen(para) * 100) // int(cap * 100)))
                for para in shape.text_frame.text.split("\n"))
    need = lines * pt * 1.45 / 72 + 0.18
    h = max(shape.height.inches, need)
    above = [sp.top.inches + sp.height.inches for sp in slide.shapes
             if sp is not shape and sp.name not in BG and sp.top is not None
             and sp.top.inches + sp.height.inches <= shape.top.inches + 0.05]
    top = min(shape.top.inches, 7.12 - h)
    if above:
        top = max(top, max(above) + 0.08)
    shape.height, shape.top = Inches(h), Inches(top)


# 別紙2 は注記が長いため、2つの表の送りを詰める
_one, _ = sh(sl[2])
_one["hcr"].top = Inches(1.40)
for _n in ("gpub", "gagr"):
    _one[_n].top = Inches(1.72)
for _c in range(7):
    _one[f"th{_c}"].top = Inches(1.99)
    _one[f"th2{_c}"].top = Inches(3.87)
for _r in range(3):
    for _c in range(7):
        _one[f"cr{_r}{_c}"].top = Inches(2.32 + _r * 0.33)
        _one[f"er{_r}{_c}"].top = Inches(4.20 + _r * 0.33)
_one["her"].top = Inches(3.30)
for _n in ("gpub2", "gagr2"):
    _one[_n].top = Inches(3.60)

for _s in sl:
    _one, _many = sh(_s)
    for _n in ("ttl", "hdr_txt"):
        for _sp in _many.get(_n, []):
            _sp.width = Inches(max(_sp.width.inches, SLIDE_W - _sp.left.inches - 0.35))
    for _n in NOTE_NAMES:
        for _sp in _many.get(_n, []):
            fit_note(_s, _sp)

for i, _s in enumerate(sl, start=1):
    _one, _ = sh(_s)
    if "pgnum" in _one:
        _one["pgnum"].width = Inches(0.85)
        set_text(_one["pgnum"], "別紙" if i == 1 else f"別紙{i - 1}")

OUT_PPTX.parent.mkdir(parents=True, exist_ok=True)
prs.save(OUT_PPTX)
print(f"saved {OUT_PPTX} ({len(sl)} slides)")
