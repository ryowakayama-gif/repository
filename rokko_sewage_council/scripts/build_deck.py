# -*- coding: utf-8 -*-
"""第2回審議会資料：全パターン（①②④）シミュレーション版の生成."""
import copy
import json
import pathlib
import re
import shutil
import tempfile
import zipfile
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches

HERE = pathlib.Path(__file__).resolve().parent.parent
BASE_PPTX = HERE / "source" / "六戸町_第2回審議会資料_案.pptx"     # 第1回時点の原稿（22ページ）
OUT_PPTX = HERE / "output" / "六戸町_第2回審議会資料.pptx"
M = json.loads((HERE / "data" / "metrics.json").read_text(encoding="utf-8"))
YI = {y: i for i, y in enumerate(
    ["R7", "R8", "R9", "R10", "R11", "R12", "R13", "R14", "R15", "R16"])}

# ---------------------------------------------------------------- 料金テーブル
BLOCKS = [("基本(0〜10㎥)", None), ("11〜20㎥", 10), ("21〜30㎥", 10), ("31〜40㎥", 10),
          ("41〜50㎥", 10), ("51〜70㎥", 20), ("71〜100㎥", 30), ("101〜150㎥", 50), ("151㎥〜", None)]
CAPS = [20, 30, 40, 50, 70, 100, 150]   # 従量区分の上限（基本10㎥の次から）

RATES = {
    "現行":   [1000, 120, 120, 130, 130, 140, 140, 140, 160],
    "①最終": [1500, 150, 155, 160, 165, 175, 185, 205, 225],
    "②最終": [1200, 180, 185, 190, 195, 205, 215, 235, 255],
    "④最終": [1400, 160, 165, 170, 175, 185, 195, 215, 235],
    # 2か年改定 R8中間（= ROUND((現行+最終)/2, 0)：モデル算定値）
    "①中間": [1250, 135, 138, 145, 148, 158, 163, 173, 193],
    "②中間": [1100, 150, 153, 160, 163, 173, 178, 188, 208],
    "④中間": [1200, 140, 143, 150, 153, 163, 168, 178, 198],
    # 3か年改定
    "①第1段": [1150, 130, 130, 140, 140, 150, 155, 160, 180],
    "①第2段": [1300, 140, 140, 150, 150, 160, 170, 180, 200],
    "②第1段": [1050, 140, 140, 150, 150, 160, 165, 170, 190],
    "②第2段": [1100, 160, 160, 170, 170, 180, 190, 200, 220],
    "④第1段": [1100, 130, 135, 140, 145, 155, 155, 165, 185],
    "④第2段": [1200, 140, 150, 150, 160, 170, 170, 190, 210],
}

PAT = {"①": "標準型", "②": "家庭軽減型", "④": "段階累進型"}


def fee(rate, vol):
    """水量 vol ㎥ の月額使用料（税抜・円）を段階累進で算定."""
    total = rate[0]
    prev = 10
    for i, cap in enumerate(CAPS, start=1):
        if vol <= prev:
            break
        total += (min(vol, cap) - prev) * rate[i]
        prev = cap
    if vol > 150:
        total += (vol - 150) * rate[8]
    return total


def tax(v):
    return int(round(v * 1.1))


def yen(v):
    return f"{v:,}円"


# ---------------------------------------------------------------- テキスト設定
def set_text(shape, text):
    """段落構造を保ったままテキストを差し替える（書式は既存ランを継承）."""
    tf = shape.text_frame
    lines = text.split("\n")
    paras = tf.paragraphs
    while len(paras) < len(lines):                       # 段落が足りなければ複製
        new = copy.deepcopy(paras[-1]._p)
        paras[-1]._p.addnext(new)
        paras = tf.paragraphs
    for i, line in enumerate(lines):
        p = paras[i]
        if not p.runs:                                   # ラン不在なら直前段落から借用
            src = next((q for q in paras if q.runs), None)
            if src is None:
                p.text = line
                continue
            p._p.append(copy.deepcopy(src.runs[0]._r))
        p.runs[0].text = line
        for r in p.runs[1:]:
            r._r.getparent().remove(r._r)
    for p in list(paras[len(lines):]):                   # 余剰段落を削除
        p._p.getparent().remove(p._p)


def sh(slide):
    """shape名 -> shape（同名は最初のもの）と、同名リストの両方を返す."""
    one, many = {}, {}
    for s in slide.shapes:
        one.setdefault(s.name, s)
        many.setdefault(s.name, []).append(s)
    return one, many


def put(slide, mapping):
    one, _ = sh(slide)
    for name, val in mapping.items():
        if name not in one:
            raise KeyError(f"shape {name!r} not found")
        set_text(one[name], val)


def put_seq(slide, name, values):
    """同名 shape が複数ある場合に順番で流し込む."""
    _, many = sh(slide)
    tgts = many[name]
    assert len(tgts) == len(values), f"{name}: {len(tgts)} shapes vs {len(values)} values"
    for t, v in zip(tgts, values):
        set_text(t, v)


def clone_shape(src_slide, src_name, dst_slide, new_name):
    """図形を複製して別（同一可）スライドへ追加する.  ID が衝突しないよう振り直す."""
    one, _ = sh(src_slide)
    el = copy.deepcopy(one[src_name]._element)
    tree = dst_slide.shapes._spTree
    used = {int(c.get("id")) for c in tree.iter()
            if c.tag.endswith("}cNvPr") and c.get("id", "").isdigit()}
    for c in el.iter():
        if c.tag.endswith("}cNvPr"):
            c.set("id", str(max(used) + 1))
            c.set("name", new_name)
    tree.append(el)
    return dst_slide.shapes[-1]


# ---------------------------------------------------------------- 指標ヘルパ
def pct(biz, metric, key, year):
    return f"{M[biz][metric][key][YI[year]] * 100:.1f}%"


def val(biz, item, year):
    return M[biz]["財政計画"][item][YI[year]]


def income(biz, key, year):
    return M[biz]["使用料収入"][key][YI[year]]


def shortfall(biz, key, year, target):
    """目標経費回収率に到達するのに不足する使用料収入（千円）と、その増収率."""
    need = val(biz, "汚水処理費", year) * target
    now = income(biz, key, year)
    return need - now, need / now - 1


def need_for_break_even(biz, year):
    """経常収支比率100%に必要な使用料収入（千円）."""
    i = YI[year]
    fp = M[biz]["財政計画"]
    return fp["経常支出"][i] - (fp["経常収入"][i] - fp["使用料収入"][i])


def ordinary_pl(biz, key, year):
    """当該パターンでの経常損益（千円）＝ 財政計画の使用料収入をパターン値に差し替えて算定."""
    i = YI[year]
    fp = M[biz]["財政計画"]
    return fp["経常収入"][i] - fp["使用料収入"][i] + income(biz, key, year) - fp["経常支出"][i]


def cum_pl(biz, key, years):
    return sum(ordinary_pl(biz, key, y) for y in years)


def unit_price(biz, key, year):
    """使用料単価（円/㎥）＝ 使用料収入 ÷ 有収水量.  交付要件の除外基準（150円/㎥）判定に使う."""
    i = YI[year]
    return M[biz]["使用料収入"][key][i] * 1000 / M[biz]["財政計画"]["有収水量"][i]


SPAN916 = ["R8", "R9", "R10", "R11", "R12", "R13", "R14", "R15", "R16"]


def signed(v):
    """会計表記の符号付き金額（千円）."""
    return f"△{-v:,.0f}" if v < 0 else f"＋{v:,.0f}"


# ---------------------------------------------------------------- 構成の組み替え
# 原稿はパターン②のみを扱う22ページ。パターン①/④用にスライドを複製し、章立てを組み直す。
DUPLICATE = [6, 6, 11, 11, 15, 15, 16, 19, 19,     # 原稿のページ番号（複製元）→ slide23〜31
             12, 16, 16, 17,                       # 3か年改定の表 → slide32〜35
             10, 8]                                # 01-1 改定しない場合 / 01-8 近隣比較 → slide36〜37
ORDER = [1, 2, 3, 4, 36, 5, 23, 6, 24, 7, 8, 37,
         9, 10, 25, 11, 26, 12, 32, 13,
         14, 27, 15, 28, 16, 33, 29, 34, 17, 35,
         18, 30, 19, 31, 20, 21, 22]

# 完成後のページ番号（本文の見出し番号 01-1 等はこれとは独立）
P = {"cover": 1, "toc": 2, "recap": 3, "sec1": 4, "nochange": 5, "eval": 6,
     "sim①": 7, "sim②": 8, "sim④": 9, "cmp": 10, "rate": 11, "peer": 12,
     "sec2": 13, "span": 14, "stg①": 15, "stg②": 16, "stg④": 17,
     "stg2": 18, "stg3": 19, "road": 20,
     "sec3": 21, "cr①": 22, "cr②": 23, "cr④": 24,
     "rec公2": 25, "rec公3": 26, "rec農2": 27, "rec農3": 28, "all2": 29, "all3": 30,
     "sec4": 31, "imp①": 32, "imp②": 33, "imp④": 34, "vol": 35, "plan": 36, "next": 37}

CT_NS = "http://schemas.openxmlformats.org/package/2006/content-types"
SLIDE_CT = ("application/vnd.openxmlformats-officedocument."
            "presentationml.slide+xml")


def duplicate_slide(root, src_no):
    """ppt/slides/slideN.xml を複製し、パッケージへの登録まで行って新しい番号を返す."""
    slides = sorted(int(m.group(1)) for m in
                    (re.fullmatch(r"slide(\d+)\.xml", f.name)
                     for f in (root / "ppt/slides").glob("slide*.xml")) if m)
    new_no = max(slides) + 1
    shutil.copy(root / f"ppt/slides/slide{src_no}.xml",
                root / f"ppt/slides/slide{new_no}.xml")
    src_rels = root / f"ppt/slides/_rels/slide{src_no}.xml.rels"
    if src_rels.exists():
        rels = src_rels.read_text(encoding="utf-8")
        rels = re.sub(r'<Relationship[^>]*notesSlide[^>]*/>', "", rels)   # ノートは引き継がない
        (root / f"ppt/slides/_rels/slide{new_no}.xml.rels").write_text(rels, encoding="utf-8")

    ct = root / "[Content_Types].xml"
    txt = ct.read_text(encoding="utf-8")
    entry = (f'<Override PartName="/ppt/slides/slide{new_no}.xml" '
             f'ContentType="{SLIDE_CT}"/>')
    ct.write_text(txt.replace("</Types>", entry + "</Types>"), encoding="utf-8")

    rp = root / "ppt/_rels/presentation.xml.rels"
    txt = rp.read_text(encoding="utf-8")
    next_id = max(int(n) for n in re.findall(r'Id="rId(\d+)"', txt)) + 1
    rid = "rId%d" % next_id
    entry = ('<Relationship Id="%s" Type="http://schemas.openxmlformats.org/'
             'officeDocument/2006/relationships/slide" '
             'Target="slides/slide%d.xml"/>' % (rid, new_no))
    rp.write_text(txt.replace("</Relationships>", entry + "</Relationships>"), encoding="utf-8")
    return new_no


def restructure(base_pptx, work_dir):
    """複製と並べ替えを済ませた中間 pptx のパスを返す."""
    root = work_dir / "unpacked"
    with zipfile.ZipFile(base_pptx) as z:
        z.extractall(root)
    for src_no in DUPLICATE:
        duplicate_slide(root, src_no)

    rels = (root / "ppt/_rels/presentation.xml.rels").read_text(encoding="utf-8")
    rid_of = {int(no): rid for rid, no in
              re.findall(r'Id="(rId\d+)"[^>]*Target="slides/slide(\d+)\.xml"', rels)}
    pres_path = root / "ppt/presentation.xml"
    pres = pres_path.read_text(encoding="utf-8")
    lst = re.search(r"<p:sldIdLst>.*?</p:sldIdLst>", pres, re.S).group(0)
    new = "<p:sldIdLst>" + "".join(
        f'<p:sldId id="{256 + i}" r:id="{rid_of[n]}"/>' for i, n in enumerate(ORDER)
    ) + "</p:sldIdLst>"
    pres_path.write_text(pres.replace(lst, new), encoding="utf-8")

    staged = work_dir / "staged.pptx"
    with zipfile.ZipFile(staged, "w", zipfile.ZIP_DEFLATED) as z:
        for f in sorted(root.rglob("*")):
            if f.is_file():
                z.write(f, f.relative_to(root).as_posix())
    return staged


_tmp = tempfile.TemporaryDirectory()
prs = Presentation(restructure(BASE_PPTX, pathlib.Path(_tmp.name)))
sl = list(prs.slides)


def S_(n):
    return sl[n - 1]


def SP(key):
    return sl[P[key] - 1]


# ================================================================ 表紙
put(SP("cover"), {
    "org_name": "六戸町 下水道使用料審議委員会",
    "title1": "第２回",
    "title2": "審議会 資料",
    "date_txt": "令和８年８月19日（水）　18:30〜",
    "loc_txt": "六戸町役場 ２階 小会議室",
})
# 表紙は元テキストより長いため、装飾（右下の四角 x=9.68"）に掛からない範囲で幅を広げる
_one, _ = sh(SP("cover"))
for _n, _w in (("org_name", 7.0), ("date_txt", 7.3), ("loc_txt", 7.3)):
    _one[_n].width = Inches(_w)

# ================================================================ 目次
put_seq(SP("toc"), "item_desc", [
    "第1回審議内容のまとめ・委員意見への対応方針",
    "改定しない場合の影響・パターン①②④の収支シミュレーション・改定率の妥当性・近隣自治体比較",
    "段階的値上げの検討（2か年vs3か年・パターン別段階単価一覧・改定ロードマップ）",
    "財政健全化目標との整合性（経常収支比率・経費回収率・目標値設定）",
    "住民への影響・説明対応（パターン別影響額・説明会計画）",
    "料金体系見直し・答申案策定・意見集約",
])

# ================================================================ 前回振り返り
put_seq(SP("recap"), "pb", [
    "公共の経費回収率47.11%・農集35.99%（令和6年度決算）と低水準。使用料改定が必要。経営持続のための抜本的対策を確認した。",
    "目標①経常収支比率100%以上、目標②経費回収率47%以上（段階的向上）を目標として設定した。",
    "パターン①②④を提示。20㎥で税抜3,000円・税込3,300円を達成。今回は3パターンすべての収支シミュレーションを提示する。",
    "5年に1回の使用料改定検証が交付要件。未達成の場合は社会資本整備総合交付金の対象外となる。",
    "2か年（R8中間→R9最終）および3か年（R8第1段→R9第2段→R10最終）の段階単価を提示・確認した。",
])

# ---------------------------------------------------------------- 章扉の説明文
put(SP("sec1"), {"desc": "本章では、まず使用料改定を行わない場合の影響を確認したうえで、"
                         "パターン①②④それぞれの収支シミュレーション結果を示し、"
                         "推奨案の選定理由と改定率の妥当性について整理します。\n"
                         "なお本章以降の指標は2か年改定を主軸として記載し、3か年改定は括弧書きで併記しています。"
                         "改定年数は第2章でご審議いただく事項です。"})
put(SP("sec3"), {"desc": "本章では、パターン①②④それぞれの経常収支比率・経費回収率の見込みを確認し、"
                         "財政健全化目標値の設定について審議いただきます。"})

# ================================================================ 01-1 改定しない場合の影響
_pl_now = {b: cum_pl(b, "現行", SPAN916) for b in ("公共", "農集")}
_pl_rev = {b: cum_pl(b, "②2", SPAN916) for b in ("公共", "農集")}
_gap = sum(_pl_rev.values()) - sum(_pl_now.values())
put(SP("nochange"), {
    "hdr_txt": "使用料改定案の検討　01-1　使用料改定を行わない場合の影響",
    "ttl": "使用料改定を行わない場合（現状維持）と改定する場合の比較",
    "ch0": "改定しない場合（現状維持）",
    "cs0": "使用料は現行のまま（基本1,000円／20㎥・税込2,420円）",
    "ph0": "✔ 当面のメリット",
    "p00": "  • 住民の月額負担は現行のまま",
    "p01": "  • 条例改正・住民説明の手続きが不要",
    "p02": "  • ただし効果は短期的なものにとどまる",
    "mh0": "✖ 想定されるリスク",
    "m00": f"  • 経常赤字が継続（R8〜R16累計 公共{signed(_pl_now['公共'])}千円・農集{signed(_pl_now['農集'])}千円）",
    "m01": "  • 一般会計繰入の増加／更新・修繕の先送り／交付金要件の未充足",
    "ch1": "改定する場合（推奨案②・2か年改定）",
    "cs1": "R8：中間単価 → R9：最終単価（20㎥・税込3,300円）",
    "ph1": "✔ 期待される効果",
    "p10": f"  • 公共はR9に経常収支比率{pct('公共','経常収支比率','②2','R9')}へ改善（単年度黒字化）",
    "p11": f"  • R8〜R16の経常損益は公共{signed(_pl_rev['公共'])}千円（改定なし比{signed(_pl_rev['公共']-_pl_now['公共'])}千円）",
    "p12": f"  • 農集も赤字幅を年6,200〜6,500千円圧縮（累計{signed(_pl_rev['農集']-_pl_now['農集'])}千円）",
    "mh1": "✖ 住民への影響",
    "m10": "  • 20㎥（3人世帯）で月額＋880円・年間＋10,560円（1日あたり約29円）",
    "m11": "  • 10㎥の少量使用世帯でも月額＋220円（パターン②の場合）",
    "rec9": ("▶ 改定しない場合、令和8〜16年度の9年間で両事業合計 約"
             f"{abs(sum(_pl_now.values()))/100000:.1f}億円の経常赤字が累積し、その補填は一般会計繰入"
             "（＝町民全体の税負担）に依存することになります。加えて社会資本整備総合交付金の交付要件"
             "（使用料改定の必要性検証・経費回収率向上ロードマップ）を満たせず、改築・更新事業の国庫補助を"
             "受けられなくなるおそれがあります。\n"
             "※ 金額は経営戦略ベースの推計値です（令和7年度決算との突合は未実施）。"
             "一般会計繰入金はブック上いずれのケースも同額を前提としているため、上表の差は経常損益の差として現れます。"),
})

# ================================================================ 01-2 評価軸
_inc = {p: income("公共", p + "2", "R9") for p in ("①", "②", "④")}
_sh9, _rt9 = shortfall("公共", "②2", "R9", 0.47)
put(SP("eval"), {
    "r00": "財政健全化への貢献度（R9使用料収入）",
    "r01": f"○ {_inc['①']:,.0f}千円", "r02": f"◎ {_inc['②']:,.0f}千円", "r03": f"○ {_inc['④']:,.0f}千円",
    "r10": "目標達成（20㎥・税抜3,000円）",
    "r11": "◎ 達成", "r12": "◎ 達成", "r13": "◎ 達成",
    "r20": "少量使用者（10㎥以下）への配慮",
    "r21": "△ 基本料+50%", "r22": "◎ 基本料+20%", "r23": "○ 基本料+40%",
    "r30": "中量使用者（20〜30㎥）の負担",
    "r31": "○ 標準", "r32": "○ 標準", "r33": "◎ 緩やか",
    "r40": "大口使用者への累進性",
    "r41": "○ 中程度", "r42": "◎ 強い", "r43": "○ 中程度",
    "r50": "総合評価", "r51": "○", "r52": "◎ 推奨", "r53": "○",
    "hdr_txt": "使用料改定案の検討　01-2　推奨パターンの絞り込み",
    "concl": (
        "▶ 推奨案：パターン②（家庭軽減型）　"
        f"【選定理由①】財政改善効果が3案中で最大（R9使用料収入{_inc['②']:,.0f}千円／公共 経費回収率"
        f"{pct('公共','経費回収率','②2','R9')}・経常収支比率{pct('公共','経常収支比率','②2','R9')}／農集 経費回収率"
        f"{pct('農集','経費回収率','②2','R9')}。いずれも3案中で最良）　"
        "【選定理由②】少量使用者（10㎥）の負担増が3案中で最小（月＋220円。①＋550円・④＋440円）\n"
        f"※ ただし公共下水道の経費回収率は目標水準47%に対しR9 {pct('公共','経費回収率','②2','R9')}で、"
        f"今回改定のみでは目標に到達しません（使用料収入ベースで{_sh9:,.0f}千円・＋{_rt9*100:.1f}%相当の不足）。"
        "R10以降の段階的適正化（次期改定：R12検証）とあわせた達成を想定しています（03-4参照）。"),
})

# ================================================================ 01-3〜5 パターン別 収支シミュレーション
SIM = {"①": "sim①", "②": "sim②", "④": "sim④"}
SUB = {"①": "01-3", "②": "01-4", "④": "01-5"}
for p, key in SIM.items():
    k2, k3 = p + "2", p + "3"
    s = SP(key)
    d = {
        "hdr_txt": f"使用料改定案の検討　{SUB[p]}　収支シミュレーション（パターン{p}採用時）",
        "ttl": f"パターン{p}（{PAT[p]}）採用時の財政収支見込み",
        "lh0": "公共下水道", "lh1": "農業集落排水",
        "g1h0": "目標① 経常収支比率（100%以上）　※（ ）内は3か年改定",
        "g1h1": "目標① 経常収支比率（100%以上）　※（ ）内は3か年改定",
        "g2h0": "目標② 経費回収率（47%以上）　※（ ）内は3か年改定",
        "g2h1": "目標② 経費回収率（47%以上）　※（ ）内は3か年改定",
    }
    for i, biz in enumerate(["公共", "農集"]):
        d[f"cr{i}0"] = f"  R7：{pct(biz,'経常収支比率','現行','R7')}"
        d[f"cr{i}1"] = f"  R8：{pct(biz,'経常収支比率',k2,'R8')}　（{pct(biz,'経常収支比率',k3,'R8')}）"
        d[f"cr{i}2"] = f"  R9：{pct(biz,'経常収支比率',k2,'R9')}　（{pct(biz,'経常収支比率',k3,'R9')}）"
        d[f"er{i}0"] = f"  R7：{pct(biz,'経費回収率','現行','R7')}"
        d[f"er{i}1"] = f"  R8：{pct(biz,'経費回収率',k2,'R8')}　（{pct(biz,'経費回収率',k3,'R8')}）"
        d[f"er{i}2"] = f"  R9：{pct(biz,'経費回収率',k2,'R9')}　（{pct(biz,'経費回収率',k3,'R9')}）"
    ko_cr10, ko_er10 = pct("公共", "経常収支比率", k2, "R10"), pct("公共", "経費回収率", k2, "R10")
    no_cr10, no_er10 = pct("農集", "経常収支比率", k2, "R10"), pct("農集", "経費回収率", k2, "R10")
    sh9, rt9 = shortfall("公共", k2, "R9", 0.47)
    d["concl6"] = (
        f"パターン{p}（{PAT[p]}）を2か年改定で採用した場合、公共下水道は経常収支比率がR9に{pct('公共','経常収支比率',k2,'R9')}"
        f"（R10以降 {ko_cr10}）となり、目標①の100%以上を達成する見込みです。経費回収率は現行34.4%から"
        f"R9 {pct('公共','経費回収率',k2,'R9')}・R10 {ko_er10}へ改善しますが、"
        f"目標②の47%には{sh9:,.0f}千円（＋{rt9*100:.1f}%）不足し、次期改定とあわせた段階的な達成となります。\n"
        f"農業集落排水は経費回収率がR9 {pct('農集','経費回収率',k2,'R9')}・R10 {no_er10}へ改善し目標②を達成する一方、"
        f"経常収支比率（目標①）はR10で{no_cr10}にとどまり、引き続き他会計繰入等による補填が必要となる見込みです。\n"
        f"3か年改定とした場合は最終単価の適用がR10となるため、目標達成が1年遅れます（上表（ ）内）。")
    put(s, d)

# ================================================================ 01-6 収支見込み比較
s = SP("cmp")
d = {"hb": "  使用料改定案の検討　01-6　収支見込み比較（パターン①②④・2か年改定）",
     "ttl": "パターン別　財政収支見込み比較（公共下水道・農業集落排水）　2か年改定ベース"}
for i, p in enumerate(["②", "④", "①"]):        # 使用料収入の大きい順に並べる
    k = p + "2"
    d[f"ph{i}"] = f"パターン{p}（{PAT[p]}）"
    d[f"ch{i}"] = "経常収支比率（目標100%以上）"
    d[f"eh{i}"] = "経費回収率（目標47%以上）"
    d[f"ach{i}"] = "経費回収率／経常収支比率"
    for j, y in enumerate(["R7", "R8", "R9"]):
        key = "現行" if y == "R7" else k
        d[f"cpr{i}{j}"] = f"  {y}：{pct('公共','経常収支比率',key,y)}"
        d[f"epr{i}{j}"] = f"  {y}：{pct('公共','経費回収率',key,y)}"
        d[f"ar{i}{j}"] = f"  {y}：{pct('農集','経費回収率',key,y)} / {pct('農集','経常収支比率',key,y)}"
d["note"] = (
    "▶ 使用料収入は②＞④＞①の順に大きく、公共下水道の経費回収率もR9で②46.1%＞④44.1%＞①43.1%。"
    "経常収支比率は3案ともR9に100%以上を達成する見込みです。農業集落排水は経費回収率が大きく改善する一方、"
    "経常収支比率は3案とも100%未満（R9で84〜86%）にとどまります。\n"
    "▶ 収入差の理由：使用料収入は基本使用料が約27%・超過（従量）使用料が約73%で、有収水量の約72%が"
    "11㎥以上の従量部分です。②は基本料金を最も低く（1,200円）抑える一方で従量単価を最も高く設定しているため、"
    "基本料金の抑制分を従量部分の増収が上回り収入が最大になります（20㎥は3案とも税抜3,000円で同額）。\n"
    f"▶ ②でも公共の経費回収率は目標47%に対しR9 {pct('公共','経費回収率','②2','R9')}で、"
    f"{_sh9:,.0f}千円（＋{_rt9*100:.1f}%）の増収が不足します（01-2・03-4参照）。")
put(s, d)

# ================================================================ 01-7 改定率の妥当性
put(SP("rate"), {
    "hdr_txt": "使用料改定案の検討　01-7　改定率の妥当性評価",
    "lh7": "青森県内自治体との比較（20㎥・税込）",
    "cn0": "六戸町（改定前）", "cp0": "2,420円", "cr0": "—",
    "cn1": "六戸町（改定後・共通）", "cp1": "3,300円", "cr1": "＋36.4%",
    "cn2": "八戸市", "cp2": "3,383円", "cr2": "—",
    "cn3": "県内平均（公共）", "cp3": "3,063円", "cr3": "—",
    "cn4": "県内平均（農集）", "cp4": "2,912円", "cr4": "—",
    "bi3": "現行比の改定率", "bv3": "＋36.4%",
    "bi4": "県内平均との差", "bv4": "＋237円",
    "eval7": "▶ 改定後の水準（税込3,300円）は青森県内平均（公共3,063円・農集2,912円）と概ね同水準。1日あたりの負担増は約29円で、住民への影響は許容範囲内と評価されます。近隣自治体との個別比較は01-8をご覧ください。",
})

# ================================================================ 01-8 近隣自治体との比較
PEER_PUB = [("おいらせ町", 3655), ("八戸市", 3383), ("七戸町", 3300),
            ("東北町", 3300), ("五戸町", 2640)]
PEER_NOU = [("黒石市（県内最高）", 4045), ("おいらせ町", 3655), ("七戸町・東北町", 3300),
            ("板柳町", 2920), ("新郷村（県内最安）", 1760)]
d = {"hdr_txt": "使用料改定案の検討　01-8　近隣自治体との使用料比較",
     "ttl": "近隣自治体との使用料比較（一般家庭 20㎥／月・税込）",
     "lh7": "公共下水道　右列＝六戸町（改定後3,300円）との差",
     "rh7": "農業集落排水（20㎥・税込）"}
for i, (name, amt) in enumerate(PEER_PUB):
    diff = amt - 3300
    d[f"cn{i}"] = name
    d[f"cp{i}"] = f"{amt:,}円"
    d[f"cr{i}"] = "±0円" if diff == 0 else (f"＋{diff:,}円" if diff > 0 else f"△{-diff:,}円")
for i, (name, amt) in enumerate(PEER_NOU):
    d[f"bi{i}"] = name
    d[f"bv{i}"] = f"{amt:,}円"
d["eval7"] = (
    "▶ 六戸町の改定後（税込3,300円）は近隣自治体の中位に位置し、七戸町・東北町と同額、"
    "八戸市（3,383円）・おいらせ町（3,655円）を下回ります。東北町は令和6年に2,640円→3,300円へ改定済みです。"
    "県内平均は公共3,063円・農集2,912円（01-7）で、改定後もこれを大きく超える水準ではありません。\n"
    "※ 文字色は六戸町の改定後3,300円との比較（赤＝高い／青＝同額／緑＝低い）。"
    "出所：第1回審議会資料「青森県内 下水道使用料比較」（令和6年3月31日現在・20㎥・税込）。"
    "三沢市は第1回資料に掲載がないため未掲載です（照会中）。")
put(SP("peer"), d)

# 金額の高低が一目で分かるよう、六戸町の改定後（3,300円）との比較で文字色を振り直す
_HEAT = {1: "B91C1C", 0: "1E3A8A", -1: "166534"}
_one, _ = sh(SP("peer"))
for _pre, _rows in (("c", PEER_PUB), ("b", PEER_NOU)):
    for _i, (_name, _amt) in enumerate(_rows):
        _rgb = RGBColor.from_string(_HEAT[(_amt > 3300) - (_amt < 3300)])
        for _n in ((f"cn{_i}", f"cp{_i}", f"cr{_i}") if _pre == "c" else (f"bi{_i}", f"bv{_i}")):
            for _p in _one[_n].text_frame.paragraphs:
                for _r in _p.runs:
                    _r.font.color.rgb = _rgb

# ================================================================ 02-1 論点整理
put(SP("span"), {
    "m00": "  • 1回の値上げ幅が大きい（基本料+100〜250円）",
    "m10": "  • 目標達成がR10まで遅延（交付金要件の充足も遅延）",
})

# ================================================================ 02-2〜4 パターン別 段階単価
STG = {"①": "stg①", "②": "stg②", "④": "stg④"}
SUB2 = {"①": "02-2", "②": "02-3", "④": "02-4"}
for p, key in STG.items():
    s = SP(key)
    lead = "推奨案（パターン②）" if p == "②" else f"パターン{p}（{PAT[p]}）"
    d = {"hdr_txt": f"段階的値上げの検討　{SUB2[p]}　段階単価一覧（パターン{p}・税抜き・円）",
         "ttl": f"{lead}の段階単価一覧"}
    cols = ["現行", p + "中間", p + "最終", p + "第1段", p + "第2段", p + "最終"]
    for r, (lab, _) in enumerate(BLOCKS):
        d[f"d{r}0"] = lab
        for c, k in enumerate(cols, start=1):
            d[f"d{r}{c}"] = f"{RATES[k][r]:,}"
    put(s, d)

# ================================================================ 02-5 段階単価比較（2か年）
s = SP("stg2")
d = {"hb": "  段階的値上げの検討　02-5　段階単価比較（パターン①②④・2か年）"}
cols = ["現行", "①中間", "①最終", "②中間", "②最終", "④中間", "④最終"]
for r, (lab, _) in enumerate(BLOCKS):
    d[f"d{r}0"] = lab
    for c, k in enumerate(cols, start=1):
        d[f"d{r}{c}"] = f"{RATES[k][r]:,}"
d["note"] = ("※ 20㎥での最終月額（税抜）は3案とも3,000円で共通。R8中間単価は算定ブックの数式どおり（現行＋最終）÷2の四捨五入。\n"
             "　 第1回資料は切り捨てで提示しており（①71〜100㎥＝162円 等）、5区分で1円の差があります。条例上の端数処理とあわせて確定が必要です。")
put(s, d)

# ================================================================ 02-6 段階単価比較（3か年）
s = SP("stg3")
d = {"hb": "  段階的値上げの検討　02-6　段階単価比較（パターン①②④・3か年）",
     "ttl": "段階単価　パターン別比較（税抜き・円）　3か年改定：R8第1段→R9第2段→R10最終",
     "gl2": "← ①標準型 3か年", "gl4": "← ②家庭軽減型 3か年", "gl6": "← ④段階累進型 3か年",
     "th0": "水量区分", "th1": "現行",
     "th2": "①R8第1段", "th3": "①R9第2段", "th4": "②R8第1段", "th5": "②R9第2段",
     "th6": "④R8第1段", "th7": "④R9第2段"}
cols3 = ["現行", "①第1段", "①第2段", "②第1段", "②第2段", "④第1段", "④第2段"]
for r, (lab, _) in enumerate(BLOCKS):
    d[f"d{r}0"] = lab
    for c, k in enumerate(cols3, start=1):
        d[f"d{r}{c}"] = f"{RATES[k][r]:,}"
d["note"] = ("※ R10の最終単価は2か年改定のR9最終単価と同一です（02-5参照）。"
             "3か年改定は各段階の値上げ幅が小さくなる一方、最終単価の適用がR10となるため財政目標の達成は1年遅れます。")
put(s, d)

# ================================================================ 02-7 ロードマップ
_road = SP("road")
put(_road, {
    "hdr_txt": "段階的値上げの検討　02-7　改定ロードマップ案",
    "ttl": "改定ロードマップ案（パターン①②④共通）",
    "tk3": "経費回収率\n目標達成確認\n効果検証",
    "tk4": "経営戦略改定\n次回改定の\n必要性検証",
})
_rn = clone_shape(SP("vol"), "note", _road, "note")
_rn.left, _rn.top = Inches(0.63), Inches(5.55)
_rn.width, _rn.height = Inches(12.08), Inches(1.42)
set_text(_rn, (
    f"▶ 今回改定の位置付け：経費回収率を34.4%→{pct('公共','経費回収率','②2','R9')}"
    "（推奨案②・2か年・R9）へ引き上げる第1段階の適正化。目標47%以上は次期改定とあわせて段階的に達成します。\n"
    "▶ 次回の検証時期：令和12年度（経営戦略の改定期）。5年に1回の使用料改定の必要性検証（交付要件）に合わせて実施します。\n"
    "▶ 再改定の判断条件：①経費回収率がロードマップの業績目標を下回る、②有収水量が推計を超えて減少する、"
    "③維持管理費・更新投資が推計を上回る場合。\n"
    "▶ 交付金への影響：要件を満たさない場合、下水道の改築・更新に係る社会資本整備総合交付金が交付対象外となります。"
    f"改定により使用料単価は{unit_price('公共','現行','R7'):.1f}円/㎥（R7）→{unit_price('公共','②2','R9'):.1f}円/㎥"
    "（②2か年・R9）となり、除外基準の150円/㎥を上回ります（交付金額・対象事業費は町に確認中）。"))

# ================================================================ 03-1〜3 パターン別 経常収支比率
CR = {"①": "cr①", "②": "cr②", "④": "cr④"}
SUB3 = {"①": "03-1", "②": "03-2", "④": "03-3"}
_need_nou = need_for_break_even("農集", "R9")
for p, key in CR.items():
    k = p + "2"
    s = SP(key)
    d = {"hdr_txt": f"財政健全化目標との整合性　{SUB3[p]}　経常収支比率の見込み（パターン{p}）",
         "ttl": f"経常収支比率の見込み（パターン{p}・{PAT[p]}／2か年改定）"}
    for i, biz in enumerate(["公共", "農集"]):
        d[f"rv{i}0"] = f"R7（現行）\n経常収支比率：{pct(biz,'経常収支比率','現行','R7')}\n（改定なし）"
        d[f"rv{i}1"] = f"R8（中間単価）\n経常収支比率：{pct(biz,'経常収支比率',k,'R8')}\n（改定開始）"
        d[f"rv{i}2"] = f"R9（最終単価）\n経常収支比率：{pct(biz,'経常収支比率',k,'R9')}\n（改定完了）"
        d[f"rv{i}3"] = f"R10〜\n経常収支比率：{pct(biz,'経常収支比率',k,'R10')}\n（安定期）"
    ok = M["公共"]["経常収支比率"][k][YI["R9"]] >= 1.0
    mult = _need_nou / income("農集", k, "R9")
    d["nt13"] = (f"✔ パターン{p}採用により、公共下水道はR9に{pct('公共','経常収支比率',k,'R9')}"
                 f"{'と目標の100%以上を達成' if ok else 'となり目標100%に届かず'}する見込みです。"
                 f"農業集落排水はR10で{pct('農集','経常収支比率',k,'R10')}にとどまり、"
                 f"100%達成には他会計繰入等による補填が引き続き必要です。\n"
                 f"※ 農業集落排水が経常収支比率100%（目標①）を単独で達成するには、R9で使用料収入 "
                 f"{_need_nou:,.0f}千円が必要です。これはパターン{p}改定後（{income('農集',k,'R9'):,.0f}千円）の"
                 f"約{mult:.1f}倍、20㎥換算で月額税込 約{3300*mult:,.0f}円に相当し、今回の改定幅で達成することは"
                 f"構造的に困難です（目標②の経費回収率47%は{pct('農集','経費回収率',k,'R9')}で達成見込み）。")
    put(s, d)

# ================================================================ 03-4〜7 経費回収率
TGT = {"R7": "—", "R8": "—", "R9": "47%以上", "R10": "50%以上", "R12": "55%以上"}
for key, biz, sub, span in [("rec公2", "公共", "03-4", "2"), ("rec公3", "公共", "03-5", "3"),
                            ("rec農2", "農集", "03-6", "2"), ("rec農3", "農集", "03-7", "3")]:
    s = SP(key)
    label = "公共下水道" if biz == "公共" else "農業集落排水"
    yrs = f"{span}か年改定"
    d = {"hdr_txt": f"財政健全化目標との整合性　{sub}　経費回収率の見込み（{label}・{yrs}）",
         "ttl": f"経費回収率の見込みと目標値（案）　{label}・{yrs}"}
    put_seq(s, "th14", ["年度", "改定なし", "パターン①", "パターン②", "パターン④", "目標水準（案）"])
    for r, y in enumerate(["R7", "R8", "R9", "R10", "R12"]):
        d[f"d14{r}0"] = y
        d[f"d14{r}1"] = pct(biz, "経費回収率", "現行", y)
        for c, p in enumerate(["①", "②", "④"], start=2):
            k = "現行" if y == "R7" else p + span
            d[f"d14{r}{c}"] = pct(biz, "経費回収率", k, y)
        d[f"d14{r}5"] = TGT[y]
    if biz == "公共" and span == "2":
        s9, r9 = shortfall("公共", "②2", "R9", 0.47)
        s10, r10 = shortfall("公共", "②2", "R10", 0.50)
        s12, r12 = shortfall("公共", "②2", "R12", 0.55)
        d["nt14"] = ("▶ 審議事項：公共下水道は今回改定のみではR9で43〜46%にとどまり、目標水準47%に到達しません。"
                     "R10以降の段階的適正化（次期改定）とあわせた目標設定について審議をお願いします。\n"
                     f"※ 推奨案②で目標に到達するのに必要な増収は、R9 47%まで{s9:,.0f}千円（＋{r9*100:.1f}%）、"
                     f"R10 50%まで{s10:,.0f}千円（＋{r10*100:.1f}%）、R12 55%まで{s12:,.0f}千円（＋{r12*100:.1f}%）です。\n"
                     "※ 経費回収率の分母＝汚水処理費（維持管理費分）＝経常費用−減価償却費−支払利息−その他営業外費用。"
                     "資本費（R7 64,832千円）は分流式下水道等に要する経費でカバーされるため全額控除しています。\n"
                     "※ R7の汚水処理費168,331千円の内訳は職員給与費5,936・動力費4,975・修繕費912・材料費1,228・"
                     "その他（主に委託料）155,280千円。R6決算の47.11%も同一基準（119,503千円）で、R7の34.4%への低下は"
                     "維持管理費の＋40.9%増によるものです（増加要因は町へ照会中・令和7年度決算と突合予定）。")
    elif biz == "公共":
        d["nt14"] = ("▶ 3か年改定では最終単価の適用がR10となるため、R9時点では38.8〜40.8%にとどまり、"
                     "R10で43.4〜46.4%に達します（2か年改定はR9で43.1〜46.1%）。いずれも目標水準47%には到達しません。\n"
                     "※ 経費回収率の分母＝汚水処理費（維持管理費分）＝経常費用−減価償却費−支払利息−その他営業外費用。"
                     "資本費（減価償却費−長期前受金戻入＋支払利息／R7 64,832千円）は分流式下水道等に要する経費でカバーされるため全額控除しています。\n"
                     "※ R10以降の水準は2か年改定と同一に収束します。")
    elif span == "2":
        d["nt14"] = ("▶ 審議事項：農業集落排水はパターン②でR9に54.2%となり目標水準47%（目標②）を達成しますが、"
                     "R11以降は有収水量の減少により低下します。次期改定（R12）での再設定を前提とした目標設定を提案します。\n"
                     f"※ 目標①（経常収支比率100%以上）は別の指標です。農業集落排水がこれを達成するにはR9で使用料収入"
                     f"{_need_nou:,.0f}千円（②改定後の約{_need_nou/income('農集','②2','R9'):.1f}倍・"
                     f"20㎥換算で月額税込 約{3300*_need_nou/income('農集','②2','R9'):,.0f}円）が必要で、"
                     "今回の改定幅では構造的に達成できません（03-2参照）。\n"
                     "※ 経費回収率の分母＝汚水処理費（維持管理費分）＝経常費用−減価償却費−支払利息−その他営業外費用。"
                     "資本費（減価償却費−長期前受金戻入＋支払利息／R7 17,236千円）は分流式下水道等に要する経費でカバーされるため全額控除しています。"
                     "R6決算の35.99%も同一基準（汚水処理費32,870千円）です。")
    else:
        d["nt14"] = ("▶ 3か年改定では最終単価の適用がR10となるため、R9時点では44.9〜48.1%となり、"
                     "目標水準47%を上回るのはパターン②のみです（2か年改定は②54.2%・④51.1%・①49.6%）。\n"
                     "※ 経費回収率の分母＝汚水処理費（維持管理費分）＝経常費用−減価償却費−支払利息−その他営業外費用。"
                     "資本費（減価償却費−長期前受金戻入＋支払利息／R7 17,236千円）は分流式下水道等に要する経費でカバーされるため全額控除しています。\n"
                     "※ R10以降の水準は2か年改定と同一に収束します。")
    put(s, d)

# ================================================================ 03-8/9 全パターン指標比較
NOTE_2 = ("▶ 使用料収入が最も大きいパターン②が両指標とも最良（R9：公共 経常101.1%・回収46.1%／農集 回収54.2%）。"
          "④・①がこれに続きます。公共下水道は3案ともR9に経常収支比率100%以上を達成する一方、"
          "農業集落排水は3案とも100%未満にとどまります。\n"
          "▶ 審議事項：農業集落排水の目標①（経常収支比率100%以上）は、今回改定では構造的に達成できないため"
          f"（R9で{_need_nou:,.0f}千円＝②改定後の約{_need_nou/income('農集','②2','R9'):.1f}倍の使用料収入が必要）、"
          "長期的な達成を目指す目標として位置付けることの可否をご審議ください。")
NOTE_3 = ("▶ 3か年改定では最終単価の適用がR10となるため、R9時点では公共の経常収支比率が98.4〜99.1%と100%に届かず、"
          "経費回収率も38.8〜40.8%にとどまります（2か年改定はR9で経常100.0〜101.1%・回収43.1〜46.1%）。\n"
          "▶ R10以降の水準は2か年改定と同一に収束します。パターン間の優劣（②＞④＞①）も2か年改定と同じです。")
for key, span, sub, note in (("all2", "2", "03-8", NOTE_2), ("all3", "3", "03-9", NOTE_3)):
    s = SP(key)
    d = {"hb": f"  財政健全化目標との整合性　{sub}　全パターン指標比較（{span}か年改定）",
         "ttl": f"パターン別　経常収支比率・経費回収率の見込み比較（{span}か年改定）"}
    for r, y in enumerate(["R7", "R8", "R9"]):
        d[f"cr{r}0"] = y
        d[f"er{r}0"] = y
        for c, (biz, p) in enumerate([(b, p) for b in ("公共", "農集") for p in ("①", "②", "④")], start=1):
            k = "現行" if y == "R7" else p + span
            d[f"cr{r}{c}"] = pct(biz, "経常収支比率", k, y)
            d[f"er{r}{c}"] = pct(biz, "経費回収率", k, y)
    d["note"] = note
    put(s, d)

# ================================================================ 04-1〜3 パターン別 影響額
IMP = {"①": "imp①", "②": "imp②", "④": "imp④"}
SUB4 = {"①": "04-1", "②": "04-2", "④": "04-3"}
HOME = [(10, "10㎥（1〜2人世帯）"), (20, "20㎥（3人世帯）　★"), (30, "30㎥（3〜4人世帯）"), (50, "50㎥（4人以上）")]
BIZC = [(50, "50㎥（小規模店舗・事務所）"), (100, "100㎥（飲食店・小売店）"),
        (200, "200㎥（宿泊施設・小規模工場）"), (500, "500㎥（学校・大型施設）")]
for p, key in IMP.items():
    s = SP(key)
    rate = RATES[p + "最終"]
    d = {"hdr_txt": f"住民への影響・説明対応　{SUB4[p]}　家庭・事業所別影響額（パターン{p}）",
         "ttl": f"家庭・事業所別　月額影響額（パターン{p}・改定完了後・税込）"}
    for pre, rows in (("hr", HOME), ("br", BIZC)):
        for r, (vol, lab) in enumerate(rows):
            cur, new = tax(fee(RATES["現行"], vol)), tax(fee(rate, vol))
            d[f"{pre}{r}0"] = lab
            d[f"{pre}{r}1"] = yen(cur)
            d[f"{pre}{r}2"] = yen(new)
            d[f"{pre}{r}3"] = f"＋{new - cur:,}円"
    up = tax(fee(rate, 20)) - tax(fee(RATES["現行"], 20))
    d["nt16"] = (f"★ 20㎥（一般家庭の標準的な使用量）での月額は税込{tax(fee(rate,20)):,}円"
                 f"（税抜{fee(rate,20):,}円）。増加額は月額＋{up:,}円、年間＋{up*12:,}円、"
                 f"1日あたり約{round(up/30.4):,}円です。※各区分の単価を段階累進で積み上げて算定。\n"
                 f"※ 改定完了後の最終単価による月額です（2か年改定はR9〜、3か年改定はR10〜。最終単価は2か年・3か年で共通）。\n"
                 f"※ 世帯人数・業種は使用水量の目安として例示したものであり、実際の使用量は世帯・事業所により異なります。")
    put(s, d)

# ================================================================ 04-4 使用水量別月額
s = SP("vol")
VOLS = [10, 15, 20, 30, 50, 100]
SHAPEROW = [0, 2, 1, 3, 4, 5]                 # 20㎥（青枠）を3行目に置くための行図形の割当て
for _r in (4, 5):                             # 行が2行足りないので複製する
    for _c in range(11):
        clone_shape(s, f"d{3 if _r == 4 else 0}{_c}", s, f"d{_r}{_c}")
d = {"hb": "  住民への影響・説明対応　04-4　使用水量別月額使用料比較（改定完了後）",
     "ttl": "使用水量別　月額使用料一覧（税込・改定完了後の最終単価）",
     "note": ("★ 青枠20㎥が一般家庭の標準使用量。改定後の差額はいずれのパターンも月額＋880円（年間＋10,560円）。税込・1か月あたり。\n"
              "※ 改定完了後の最終単価による月額です（2か年改定はR9〜、3か年改定はR10〜。最終単価は2か年・3か年で共通）。\n"
              "※ 使用量の目安：10㎥＝1〜2人世帯、15〜20㎥＝2〜3人世帯、30㎥＝3〜4人世帯、50㎥＝4人以上の世帯・小規模事業所、100㎥＝飲食店・小売店等（あくまで目安）。")}
for i, vol in enumerate(VOLS):
    r = SHAPEROW[i]
    cur = tax(fee(RATES["現行"], vol))
    d[f"d{r}0"] = f"{vol}㎥"
    d[f"d{r}1"] = yen(cur)
    for c, p in enumerate(["①", "②", "④"]):
        v = fee(RATES[p + "最終"], vol)
        base = 2 + c * 3
        d[f"d{r}{base}"] = yen(v)
        d[f"d{r}{base+1}"] = yen(tax(v))
        d[f"d{r}{base+2}"] = f"＋{tax(v)-cur:,}円"
put(s, d)

# ================================================================ 04-5 住民説明会
put(SP("plan"), {"hdr_txt": "住民への影響・説明対応　04-5　住民説明会の実施計画（案）"})

# ================================================================ レイアウト調整
# 差し替えたテキストが元より長い箇所があるため、見出し・注記の器を実寸に合わせて調整する
SLIDE_W = 13.333
NOTE_NAMES = ("nt13", "nt14", "nt16", "note", "eval7", "concl", "rec9", "concl6")
BG = ("bg1", "bg2", "hdr_bg", "hdr_line", "hdr_accent", "hb", "hl", "ha")


def _vlen(s):
    """全角換算の文字数（半角は0.5）."""
    return sum(0.5 if ord(c) < 0x2000 else 1.0 for c in s)


def _fontpt(shape, default=11.0):
    for para in shape.text_frame.paragraphs:
        for r in para.runs:
            if r.font.size:
                return r.font.size.pt
    return default


def fit_note(slide, shape):
    """注記の高さを行数に合わせ、上の要素とページ番号の間に収める."""
    pt = _fontpt(shape)
    cap = max(1.0, (shape.width.inches - 0.35) * 72 / pt)
    lines = sum(max(1, -(-int(_vlen(para) * 100) // int(cap * 100)))
                for para in shape.text_frame.text.split("\n"))
    need = lines * pt * 1.45 / 72 + 0.18
    h = max(shape.height.inches, need)
    above = [sp.top.inches + sp.height.inches for sp in slide.shapes
             if sp is not shape and sp.name not in BG and sp.top is not None
             and sp.top.inches + sp.height.inches <= shape.top.inches + 0.05]
    top = min(shape.top.inches, 7.12 - h)
    if above:
        top = max(top, max(above) + 0.08)
    shape.height, shape.top = Inches(h), Inches(top)


# 01-2 は結論欄が5行になるため、評価表の送りを詰める
_one, _ = sh(SP("eval"))
for _r in range(6):
    for _c in range(4):
        _sp = _one[f"r{_r}{_c}"]
        _sp.top, _sp.height = Inches(2.52 + _r * 0.50), Inches(0.46)

# 01-6 は注記に収入差の説明を加えたため、3カラムの指標ブロック全体を上へ詰める
_STACK = [("ph", 1.50, 0.33), ("pub", 1.85, 0.25), ("ch", 2.13, 0.24),
          ("cpr0", 2.40, 0.26), ("cpr1", 2.68, 0.26), ("cpr2", 2.96, 0.26),
          ("eh", 3.25, 0.24),
          ("epr0", 3.52, 0.26), ("epr1", 3.80, 0.26), ("epr2", 4.08, 0.26),
          ("agr", 4.37, 0.25), ("ach", 4.65, 0.24),
          ("ar0", 4.92, 0.26), ("ar1", 5.20, 0.26), ("ar2", 5.48, 0.26)]
_one, _ = sh(SP("cmp"))
for _i in range(3):
    for _pre, _t, _h in _STACK:
        _n = f"{_pre[:-1]}{_i}{_pre[-1]}" if _pre[-1].isdigit() else f"{_pre}{_i}"
        _sp = _one[_n]
        _sp.top, _sp.height = Inches(_t), Inches(_h)

# 03-4〜7 は注記が長いため送りを詰め、あわせて列の文字色を意味づけに合わせる
# （原稿は列ごとに交互配色で、推奨案②が「赤＝未達」に見えてしまうため）
_COL = {1: ("B91C1C", False), 2: ("334155", False), 3: ("166534", True),
        4: ("334155", False), 5: ("1E3A8A", True)}
for _k in ("rec公2", "rec公3", "rec農2", "rec農3"):
    _one, _ = sh(SP(_k))
    for _r in range(5):
        for _c in range(6):
            _sp = _one[f"d14{_r}{_c}"]
            _sp.top, _sp.height = Inches(2.13 + _r * 0.62), Inches(0.58)
            if _c in _COL:
                _rgb, _bold = _COL[_c]
                for _p in _sp.text_frame.paragraphs:
                    for _run in _p.runs:
                        _run.font.color.rgb = RGBColor.from_string(_rgb)
                        _run.font.bold = _bold

# 03-1〜3 は注記に農集の必要収入を追記したため、経常収支比率カードの送りを詰める
for _k in ("cr①", "cr②", "cr④"):
    _one, _ = sh(SP(_k))
    for _i in range(2):
        for _r in range(4):
            _sp = _one[f"rv{_i}{_r}"]
            _sp.top, _sp.height = Inches(2.47 + _r * 0.86), Inches(0.82)

# 02-5・02-6 は元から最終行と注記が重なるため、データ行の送りを詰める
for _k in ("stg2", "stg3"):
    _one, _ = sh(SP(_k))
    for _r in range(9):
        for _c in range(8):
            _sp = _one[f"d{_r}{_c}"]
            _sp.top, _sp.height = Inches(2.13 + _r * 0.47), Inches(0.44)

# 01-6 は②④①の順に並べ替えたため、パターン見出しの色を①青・②緑・④橙の凡例に合わせ直す
_one, _ = sh(SP("cmp"))
for _n, _rgb in (("ph0", "166534"), ("ph1", "B45309"), ("ph2", "1E3A8A")):   # ②緑・④橙・①青
    _one[_n].fill.solid()
    _one[_n].fill.fore_color.rgb = RGBColor.from_string(_rgb)

# 04-4 は6行に増やしたため、見出し・データ行の送りを組み直す
_one, _many = sh(SP("vol"))
for _n in ("gl2", "gl5", "gl8"):
    _one[_n].top = Inches(1.44)
for _c in range(11):
    _sp = _one[f"th{_c}"]
    _sp.top, _sp.height = Inches(1.70), Inches(0.30)
for _i in range(6):
    _r = SHAPEROW[_i]
    for _c in range(11):
        _sp = _one[f"d{_r}{_c}"]
        _sp.top, _sp.height = Inches(1.96 + _i * 0.71), Inches(0.65)
        if _i != 2:                                  # 20㎥（青枠）以外はゼブラを振り直す
            _sp.fill.solid()
            _sp.fill.fore_color.rgb = RGBColor.from_string("F8FAFF" if _i % 2 == 0 else "FFFFFF")

# 事業所テーブルの「差額」列が元から狭く数値が折り返すため、列幅を配分し直す（04-1〜3）
for _k in ("imp①", "imp②", "imp④"):
    _one, _ = sh(SP(_k))
    _one["bh"].width = Inches(5.60)
    for _r in range(4):
        for _c, (_l, _w) in enumerate([(7.05, 2.20), (9.35, 1.05), (10.50, 1.15), (11.75, 1.15)]):
            _sp = _one[f"br{_r}{_c}"]
            _sp.left, _sp.width = Inches(_l), Inches(_w)

# 段階単価表の見出しは2行（例「R9 最終／（2か年）」）で元の行高に収まらないため、
# 見出し行を高くし、区分バーを少し上へ逃がす（02-2〜4）
for _k in ("stg①", "stg②", "stg④"):
    _one, _many = sh(SP(_k))
    for _n in ("g2yr", "g3yr"):
        _one[_n].top = Inches(1.68)
    for _sp in _many["th"]:
        _sp.top, _sp.height = Inches(1.97), Inches(0.45)

# 右側テーブルがスライド右端をはみ出すため内側へ寄せる（01-7・01-8）
for _k in ("rate", "peer"):
    _one, _ = sh(SP(_k))
    for _r in range(5):
        _one[f"bi{_r}"].width = Inches(3.95)
        _one[f"bv{_r}"].left, _one[f"bv{_r}"].width = Inches(11.13), Inches(2.13)

# 01-8 は3列目に差額（例「△660円」）が入り既定幅では折り返すため、左表の列幅を配分し直す
_one, _ = sh(SP("peer"))
for _r in range(5):
    for _n, _l, _w in ((f"cn{_r}", 0.63, 2.45), (f"cp{_r}", 3.20, 1.35), (f"cr{_r}", 4.63, 1.28)):
        _one[_n].left, _one[_n].width = Inches(_l), Inches(_w)

for _s in sl:
    _one, _many = sh(_s)
    for _n in ("ttl", "hdr_txt"):                       # 見出しは右方向に余白があるので広げる
        for _sp in _many.get(_n, []):
            _sp.width = Inches(max(_sp.width.inches, SLIDE_W - _sp.left.inches - 0.35))
    for _n in NOTE_NAMES:
        for _sp in _many.get(_n, []):
            fit_note(_s, _sp)

# ================================================================ ページ番号
for i, s in enumerate(sl, start=1):
    one, _ = sh(s)
    if "pgnum" in one:
        set_text(one["pgnum"], str(i))

OUT_PPTX.parent.mkdir(parents=True, exist_ok=True)
prs.save(OUT_PPTX)
print(f"saved {OUT_PPTX} ({len(sl)} slides)")
