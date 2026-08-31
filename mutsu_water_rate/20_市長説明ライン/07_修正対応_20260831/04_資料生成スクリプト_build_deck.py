# -*- coding: utf-8 -*-
import copy, warnings
warnings.filterwarnings('ignore')
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

SP="/tmp/claude-0/-home-user-repository/05ae8dfb-f2f7-5c38-805c-45df053cf245/scratchpad"
FONT="BIZ UDPゴシック"
NAVY=RGBColor(0x1C,0x3C,0x60); RED=RGBColor(0xBE,0x41,0x41); INK=RGBColor(0x19,0x1C,0x23)
WHITE=RGBColor(0xFF,0xFF,0xFF); HDRBG=RGBColor(0x1C,0x3C,0x60)
LBLBG=RGBColor(0xE1,0xEB,0xF7); ROW_A=RGBColor(0xFF,0xFF,0xFF); ROW_B=RGBColor(0xF4,0xF7,0xFB)

# ---- 算定値（recalc.py の検証済みロジック） -------------------------------
exec(open(f"{SP}/recalc.py").read().split("YRS=")[0])
A=run(B6_ORIG, False)   # 改定あり（納品モデルのまま）
Bn=run(0.0,     False)  # 改定なし
def m(v):
    n=round(v/1000.0)
    return f"△{abs(n):,}" if n<0 else f"{n:,}"
def pc(v): return f"{v:.1f}%"
IDX=range(1,10)                       # R8..R16

def setcell(cell, text, *, size=8.5, bold=False, color=INK, bg=None, align=PP_ALIGN.CENTER):
    tf=cell.text_frame
    tf.clear()
    p=tf.paragraphs[0]; p.alignment=align
    r=p.add_run(); r.text=text
    r.font.size=Pt(size); r.font.bold=bold; r.font.name=FONT
    r.font.color.rgb=color
    # 東アジア字体も明示
    from pptx.oxml.ns import qn
    rPr=r._r.get_or_add_rPr()
    for tag in ("a:ea","a:cs"):
        e=rPr.find(qn(tag))
        if e is None:
            e=rPr.makeelement(qn(tag),{}); rPr.append(e)
        e.set("typeface",FONT)
    if bg is not None:
        cell.fill.solid(); cell.fill.fore_color.rgb=bg

prs=Presentation(f"{SP}/base_deck.pptx")

# =========================================================================
# 1) P9：「補填財源残高（改定なし）」R11〜R16 を補完し、書式を整える
# =========================================================================
s9=prs.slides[8]
t9=[sh for sh in s9.shapes if getattr(sh,'has_table',False)][0].table
vals_nashi=[m(Bn["r72"][i]) for i in IDX]          # R8..R16
setcell(t9.cell(7,0),"補填財源残高（改定なし）",bold=True,color=RED,bg=LBLBG,align=PP_ALIGN.LEFT)
for j,v in enumerate(vals_nashi):
    setcell(t9.cell(7,j+1), v, color=RED, bg=ROW_A)
# 「なし=赤／あり=紺」で新P10と配色を統一
setcell(t9.cell(8,0),"補填財源残高（改定あり）",bold=True,color=NAVY,bg=LBLBG,align=PP_ALIGN.LEFT)
for j,v in enumerate([m(A["r72"][i]) for i in IDX]):
    setcell(t9.cell(8,j+1), v, color=NAVY, bg=ROW_B)

# P9脚注：新スライド挿入で「スライド10」→「スライド11」、2行の説明を追記
for sh in s9.shapes:
    if sh.has_text_frame and "起点とした簡便推計" in sh.text_frame.text:
        tf=sh.text_frame; p=tf.paragraphs[0]
        for r in list(p.runs)[1:]: r._r.getparent().remove(r._r)
        p.runs[0].text=("単位：百万円。補填財源残高はR7末183百万円を起点とした簡便推計。"
                        "「改定なし」は現行料金を維持した場合、「改定あり」はR11以降にスライド11の"
                        "資産維持費0％ケース（必要改定率7.3％）を反映した場合。")

# =========================================================================
# 2) 新スライドを P9 の複製として作成し、P10 に挿入
# =========================================================================
def duplicate(prs, idx):
    src=prs.slides[idx]
    dst=prs.slides.add_slide(src.slide_layout)
    for shp in list(dst.shapes): shp._element.getparent().remove(shp._element)
    for shp in src.shapes: dst.shapes._spTree.append(copy.deepcopy(shp._element))
    return dst

new=duplicate(prs,8)
shapes=list(new.shapes)
tbl=[sh for sh in shapes if getattr(sh,'has_table',False)][0].table

ROWS=[("補填財源残高","改定なし", [m(Bn["r72"][i]) for i in IDX], RED),
      ("補填財源残高","改定あり", [m(A ["r72"][i]) for i in IDX], NAVY),
      ("現預金残高",  "改定なし", [m(Bn["r59"][i]) for i in IDX], RED),
      ("現預金残高",  "改定あり", [m(A ["r59"][i]) for i in IDX], NAVY),
      ("当年度純損益","改定なし", [m(Bn["r28"][i]) for i in IDX], RED),
      ("当年度純損益","改定あり", [m(A ["r28"][i]) for i in IDX], NAVY),
      ("経常収支比率","改定なし", [pc(Bn["r31"][i]) for i in IDX], RED),
      ("経常収支比率","改定あり", [pc(A ["r31"][i]) for i in IDX], NAVY)]

setcell(tbl.cell(0,0),"区分",size=8.5,bold=True,color=WHITE,bg=HDRBG,align=PP_ALIGN.LEFT)
for j,y in enumerate(["R8","R9","R10","R11","R12","R13","R14","R15","R16"]):
    setcell(tbl.cell(0,j+1),y,bold=True,color=WHITE,bg=HDRBG)
for i,(ind,case,vals,col) in enumerate(ROWS):
    rbg = ROW_A if (i//2)%2==0 else ROW_B
    setcell(tbl.cell(i+1,0), f"{ind}　{case}", bold=True, color=col, bg=LBLBG, align=PP_ALIGN.LEFT)
    for j,v in enumerate(vals):
        setcell(tbl.cell(i+1,j+1), v, color=col, bg=rbg)

TXT={0:"料金改定を行わない場合の影響",
     1:"現行料金のままでは、補填財源・現預金はR16まで一度も回復しない",
     2:"10"}
for i,t in TXT.items():
    tf=shapes[i].text_frame; p=tf.paragraphs[0]
    for r in list(p.runs)[1:]: r._r.getparent().remove(r._r)
    p.runs[0].text=t

body=("・料金改定を行わない場合、補填財源残高はR9にマイナスへ転じた後、R16末△123百万円まで回復しない。\n"
      "・現預金もR10以降マイナスで推移し、工事代金・企業債償還の支払財源を確保できない。\n"
      "・純損益はR9に赤字転落し、R16には△126百万円まで拡大する。改定なしでは経常収支比率も100％を下回り続ける。\n"
      "・資産維持費0％（必要改定率7.3％）で改定した場合との差は、R16末で補填財源・現預金とも約1,006百万円。")
tf=shapes[4].text_frame; tf.clear()
for k,line in enumerate(body.split("\n")):
    p=tf.paragraphs[0] if k==0 else tf.add_paragraph()
    r=p.add_run(); r.text=line
    r.font.size=Pt(11); r.font.name=FONT; r.font.color.rgb=INK
    from pptx.oxml.ns import qn
    rPr=r._r.get_or_add_rPr()
    for tag in ("a:ea","a:cs"):
        e=rPr.find(qn(tag))
        if e is None: e=rPr.makeelement(qn(tag),{}); rPr.append(e)
        e.set("typeface",FONT)

tf=shapes[5].text_frame; p=tf.paragraphs[0]
for r in list(p.runs)[1:]: r._r.getparent().remove(r._r)
p.runs[0].text=("単位：百万円（経常収支比率は％）。「改定なし」は現行料金を維持した場合、"
                "「改定あり」はR11以降に必要改定率7.3％を反映した場合。"
                "投資計画・企業債条件は両ケース共通のため、資本的収支不足額は同額。"
                "補填財源はR7末183百万円を起点とした簡便推計。")

# P10 の位置へ移動
lst=prs.slides._sldIdLst; ids=list(lst)
lst.remove(ids[-1]); lst.insert(9, ids[-1])

# =========================================================================
# 3) 以降のスライドのページ番号を繰り下げ
# =========================================================================
for i in range(10, len(prs.slides._sldIdLst)):
    sl=prs.slides[i]
    for sh in sl.shapes:
        if sh.has_text_frame:
            txt=sh.text_frame.text.strip()
            if txt.isdigit() and int(txt)==i:      # 旧番号 == 新index
                p=sh.text_frame.paragraphs[0]
                for r in list(p.runs)[1:]: r._r.getparent().remove(r._r)
                p.runs[0].text=str(i+1)

out=f"{SP}/mutsu_mayor_briefing_v9_hoten_kaiteinashi.pptx"
prs.save(out); print("保存:",out)
