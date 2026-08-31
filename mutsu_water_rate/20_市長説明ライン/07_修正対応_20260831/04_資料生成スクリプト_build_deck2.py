# -*- coding: utf-8 -*-
import warnings; warnings.filterwarnings('ignore')
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

SP="/tmp/claude-0/-home-user-repository/05ae8dfb-f2f7-5c38-805c-45df053cf245/scratchpad"
FONT="BIZ UDPゴシック"
RED=RGBColor(0xBE,0x41,0x41)
LBLBG=RGBColor(0xE1,0xEB,0xF7)      # 見出し列の薄い青
WHITE=RGBColor(0xFF,0xFF,0xFF)      # 白
PALE =RGBColor(0xF4,0xF7,0xFB)      # 薄い青（交互用）

exec(open(f"{SP}/recalc.py").read().split("YRS=")[0])
A =run(B6_ORIG, False)   # 改定あり（現行モデル）
Bn=run(0.0,     False)   # 改定なし
def m(v):
    n=round(v/1000.0)
    return f"△{abs(n):,}" if n<0 else f"{n:,}"
IDX=range(1,10)          # R8..R16

def setcell(cell,text,*,size=8.5,bold=False,color=RED,bg=None,align=PP_ALIGN.CENTER):
    tf=cell.text_frame; tf.clear()
    p=tf.paragraphs[0]; p.alignment=align
    r=p.add_run(); r.text=text
    r.font.size=Pt(size); r.font.bold=bold; r.font.name=FONT; r.font.color.rgb=color
    rPr=r._r.get_or_add_rPr()
    for tag in ("a:ea","a:cs"):
        e=rPr.find(qn(tag))
        if e is None: e=rPr.makeelement(qn(tag),{}); rPr.append(e)
        e.set("typeface",FONT)
    if bg is not None:
        cell.fill.solid(); cell.fill.fore_color.rgb=bg

# 市から受領した原本（15ページ）を起点にする＝追加ページなし・ページ番号は元のまま
prs=Presentation(f"{SP}/base_deck.pptx")
t9=[sh for sh in prs.slides[8].shapes if getattr(sh,'has_table',False)][0].table

# 表の既存パターン（r1白 / r2薄青 …／r6薄青）を継承 → r7=白、r8=薄青
for ri,(label,vals,rowbg) in enumerate([
        ("補填財源残高（改定なし）",[m(Bn["r72"][i]) for i in IDX], WHITE),
        ("補填財源残高（改定あり）",[m(A ["r72"][i]) for i in IDX], PALE)], start=7):
    setcell(t9.cell(ri,0), label, bold=True, color=RED, bg=LBLBG, align=PP_ALIGN.LEFT)
    for j,v in enumerate(vals):
        setcell(t9.cell(ri,j+1), v, color=RED, bg=rowbg)

# 脚注：2行の意味を明記（スライド参照は原本どおり「スライド10」で正しい）
for sh in prs.slides[8].shapes:
    if sh.has_text_frame and "起点とした簡便推計" in sh.text_frame.text:
        p=sh.text_frame.paragraphs[0]
        for r in list(p.runs)[1:]: r._r.getparent().remove(r._r)
        p.runs[0].text=("単位：百万円。補填財源残高はR7末183百万円を起点とした簡便推計。"
                        "「改定なし」は現行料金を維持した場合、「改定あり」はR11以降にスライド10の"
                        "資産維持費0％ケース（必要改定率7.3％）を反映した場合。")

out=f"{SP}/mutsu_mayor_briefing_v9r2.pptx"
prs.save(out); print("保存:",out)
